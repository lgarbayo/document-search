"""
services/document_extractor.py — Pipeline de Ingesta y Procesamiento de Documentos.

Este módulo es el corazón de la extracción de conocimiento de Meiga. Soporta una 
amplia variedad de formatos corporativos, transformando archivos binarios en 
texto limpio y estructurado listo para ser vectorizado.

Capacidades principales:
    - Extracción Multi-formato: PDF (PyMuPDF + OCR), Office (Word, PPT, Excel), 
      Imágenes (Tesseract), Datos (JSON, XML, CSV) y Marcado (MD, HTML).
    - Metadatos Avanzados: Integración con ExifTool para extraer autoría, software 
      creador y fechas técnicas.
    - Inteligencia de Categorización: Heurísticas basadas en palabras clave y 
      nombres de archivo para auto-clasificar documentos.
    - Limpieza y Fragmentación: Normalización de texto y "Semantic Chunking" 
      respetando límites de oraciones.

Pipeline de Procesamiento:
    extract_document_content() ➔ clean_text() ➔ chunk_text() ➔ deduplicate_chunks()
"""

import csv
import hashlib
import io
import os
import re
from pathlib import Path

import fitz  # PyMuPDF
import logging

try:
    from pdf2image import convert_from_path
    from PIL import Image, ImageOps
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

# Extensiones soportadas
SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".csv", ".xlsx", ".png", ".jpg", ".jpeg",
                        ".docx", ".pptx", ".json", ".xml", ".md", ".html"}


def extract_document_content(file_path: str) -> tuple[str, dict]:
    """
    Coordina la extracción de texto y la recolección de metadatos de un archivo.

    Es el punto de entrada principal para procesar cualquier archivo individual. 
    Detecta el tipo de archivo por su extensión y delega la extracción al 
    sub-extractor correspondiente.

    Args:
        file_path (str): Ruta absoluta al archivo en el sistema de archivos del servidor.

    Returns:
        tuple[str, dict]: Una tupla conteniendo:
            - texto_crudo (str): El contenido textual completo del documento.
            - metadatos (dict): Diccionario con información técnica (tamaño, autor, 
              categoría, datos ExifTool).

    Raises:
        FileNotFoundError: Si el archivo no existe en la ruta proporcionada.
        ValueError: Si el formato de archivo no es soportado o no se pudo extraer texto.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Archivo no encontrado: {file_path}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Formato no soportado: {ext}. "
            f"Formatos válidos: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    # Dispatch según extensión
    extractors = {
        ".pdf": _extract_pdf,
        ".txt": _extract_txt,
        ".csv": _extract_csv,
        ".xlsx": _extract_xlsx,
        ".png": _extract_img,
        ".jpg": _extract_img,
        ".jpeg": _extract_img,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".json": _extract_json,
        ".xml": _extract_xml,
        ".md": _extract_txt,
        ".html": _extract_html,
    }

    text, extra_meta = extractors[ext](file_path)

    if not text.strip():
        raise ValueError(f"El archivo no contiene texto extraíble: {file_path}")

    # Extraer metadatos
    file_size = os.path.getsize(file_path)
    category = _infer_category(text, ext, filename=path.stem)

    metadata = {
        "extension": ext,
        "file_size_bytes": file_size,
        "category": category,
    }
    
    # Extraer metadatos técnicos avanzados con ExifTool
    try:
        import subprocess
        import json
        from datetime import datetime
        result = subprocess.run(
            ["exiftool", "-j", str(path)],
            capture_output=True,
            text=True,
            check=False
        )
        if result.stdout:
            parsed_exif = json.loads(result.stdout)
            if isinstance(parsed_exif, list) and len(parsed_exif) > 0:
                exif_data = parsed_exif[0]
                for k in ["Directory", "SourceFile", "FileName", "ExifToolVersion"]:
                    exif_data.pop(k, None)
                metadata["exif_metadata"] = exif_data

                # Intentar extraer el año (YYYY) de CreateDate o FileModifyDate para filtrado numérico en Qdrant
                date_str = exif_data.get("CreateDate") or exif_data.get("FileModifyDate")
                if date_str and isinstance(date_str, str) and len(date_str) >= 4:
                    try:
                        metadata["exif_year"] = int(date_str[:4])
                        # Intentar extraer mes (MM) también
                        if len(date_str) >= 7:
                            try:
                                month = int(date_str[5:7])
                                if 1 <= month <= 12:
                                    metadata["exif_month"] = month
                            except (ValueError, IndexError):
                                pass
                    except ValueError:
                        pass

                # Extraer y almacenar explícitamente Author
                author = exif_data.get("Author")
                if author and isinstance(author, str) and author.strip():
                    metadata["author"] = author.strip()

                # Extraer y almacenar explícitamente Creator (ej., software que creó el PDF)
                creator = exif_data.get("Creator")
                if creator and isinstance(creator, str) and creator.strip():
                    metadata["creator"] = creator.strip()
    except Exception as e:
        logger.warning(f"No se pudo extraer metadata ExifTool para {file_path}: {e}")
    
    # Combinar con los metadatos específicos del extractor (ej. autor de PDF)
    if extra_meta:
        metadata.update(extra_meta)

    return text, metadata


def _infer_category(text: str, ext: str, filename: str = "") -> str:
    """
    Aplica heurísticas lingüísticas para clasificar el documento automáticamente.

    Combina el análisis de palabras clave en el contenido con el análisis del 
    nombre del archivo. Se asignan pesos mayores (3x) a las coincidencias en el 
    nombre del archivo, ya que suelen ser más indicativas del propósito del documento.

    Categorías soportadas: RRHH, Finanzas, Legal, Técnico, Comercial, 
    Sostenibilidad, IT/Soporte y General.

    Args:
        text (str): El texto extraído (se analizan los primeros 8000 caracteres).
        ext (str): Extensión del archivo para sesgos de formato.
        filename (str): Nombre del archivo (sin ruta) para análisis de contexto.

    Returns:
        str: Nombre de la categoría más probable.
    """
    text_lower = text[:8000].lower()
    filename_lower = filename.lower() if filename else ""

    categories = {
        "RRHH": [
            "vacaciones", "nómina", "nóminas", "contrato laboral", "empleado",
            "salario", "despido", "baja laboral", "ausencia", "beneficios",
            "recursos humanos", "selección", "candidato", "onboarding",
            "teletrabajo", "jornada", "convenio", "antigüedad", "plantilla",
            "offer letter", "puesto", "retribución", "compensación",
        ],
        "Finanzas": [
            "factura", "facturación", "iva", "balance", "gastos",
            "presupuesto", "ingresos", "contabilidad", "fiscal", "pago",
            "pedido", "proveedor", "cobro", "deuda", "amortización",
            "cuenta", "financiero", "trimestre", "resultado", "margen",
            "ventas", "revenue", "coste", "precio", "descuento", "euro",
        ],
        "Legal": [
            "acuerdo", "confidencialidad", "ley", "decreto", "demanda",
            "litigio", "cláusula", "poder notarial", "notario", "nda",
            "licitación", "pliego", "condiciones", "adjudicación",
            "constitución", "estatutos", "escritura", "mercantil",
            "propiedad intelectual", "rgpd", "protección de datos",
        ],
        "Técnico": [
            "api", "servidor", "base de datos", "código", "despliegue",
            "arquitectura", "software", "hardware", "firmware",
            "iot", "sensor", "panel", "manual", "especificación",
            "ficha técnica", "configuración", "instalación", "versión",
        ],
        "Comercial": [
            "propuesta", "oferta", "cliente", "comercial", "catálogo",
            "producto", "servicio", "proyecto", "licitación", "smart",
            "solución", "demo", "piloto", "partnership", "colaboración",
        ],
        "Sostenibilidad": [
            "sostenibilidad", "medioambiental", "emisiones", "carbono",
            "residuos", "reciclaje", "ods", "rsc", "impacto ambiental",
            "huella", "certificación ambiental", "iso 14001", "energía",
        ],
        "IT/Soporte": [
            "incidencia", "soporte", "ticket", "resolución", "sla",
            "equipo", "inventario", "activo", "licencia", "renovación",
            "backup", "antivirus", "red", "vpn", "usuario",
        ],
    }

    scores = {cat: 0 for cat in categories}

    for cat, keywords in categories.items():
        for kw in keywords:
            # Peso 1× por cada aparición en el texto
            scores[cat] += text_lower.count(kw)
            # Peso 3× si la keyword aparece en el nombre del archivo
            if kw in filename_lower:
                scores[cat] += 3

    # Extra: si es CSV/XLSX sin coincidencias claras, sesgo a Finanzas o IT/Soporte
    if ext in [".csv", ".xlsx"]:
        if scores["Finanzas"] == 0 and scores["IT/Soporte"] == 0:
            scores["Finanzas"] += 1

    best_category = max(scores, key=scores.get)

    if scores[best_category] == 0:
        return "General"

    return best_category


def _extract_pdf(file_path: str) -> tuple[str, dict]:
    """
    Extrae texto de documentos PDF con una estrategia híbrida robusta.

    1. Intento Digital (PyMuPDF): Extrae capas de texto digital, siendo 
       el método más rápido y preciso.
    2. Fallback de OCR (Tesseract): Si el PDF está escaneado o tiene 
       muy poco texto (como faxes o imágenes), se activa el reconocimiento 
       óptico de caracteres página a página.

    Args:
        file_path (str): Ruta al PDF.

    Returns:
        tuple: (texto_extraído, metadatos_nativos_pdf).
    """
    try:
        doc = fitz.open(file_path)
    except Exception as e:
        raise ValueError(f"No se pudo abrir el PDF: {e}")

    # Extraer metadatos nativos del PDF
    pdf_meta = {}
    if doc.metadata:
        for key in ["author", "creator", "title", "subject", "keywords", "producer"]:
            val = doc.metadata.get(key)
            if val and isinstance(val, str) and val.strip():
                # Normalizar nombres de metadatos (ej., creator -> empresa/software)
                pdf_meta[key.lower()] = val.strip()

    text_parts = []
    for page in doc:
        page_text = page.get_text("text")
        if page_text.strip():
            text_parts.append(page_text)

    doc.close()
    extracted_text = "\n".join(text_parts)

    # Heurística: Si el PDF tiene menos de 50 caracteres (ej. una imagen escaneada gigante),
    # intentamos usar Tesseract OCR como fallback.
    if len(extracted_text.strip()) < 50 and OCR_AVAILABLE:
        logger.info(f"Poco texto extraído de {file_path}. Iniciando Tesseract OCR fallback...")
        try:
            # Convertimos las páginas del PDF a imágenes (máximo 300 DPI)
            images = convert_from_path(file_path, dpi=300)
            ocr_text_parts = []
            for img in images:
                # Corregir rotación EXIF si la hay
                img = ImageOps.exif_transpose(img)
                # Extraemos texto usando los modelos en español e inglés
                page_ocr = pytesseract.image_to_string(img, lang="spa+eng", config="--psm 1")
                ocr_text_parts.append(page_ocr)
            
            extracted_text = "\n".join(ocr_text_parts)
            logger.info("OCR finalizado exitosamente.")
        except Exception as e:
            logger.warning(f"Fallo en el OCR fallback para {file_path}: {e}")

    return extracted_text, pdf_meta


def _extract_img(file_path: str) -> tuple[str, dict]:
    """
    Extrae texto de archivos de imagen (PNG, JPG, JPEG) mediante OCR.

    Utiliza Tesseract OCR con modelos de lenguaje dual (español e inglés) 
    para maximizar la precisión en documentos técnicos bilingües.

    Args:
        file_path (str): Ruta a la imagen.

    Returns:
        tuple[str, dict]: Texto extraído y diccionario de metadatos (vacío por defecto).

    Raises:
        ValueError: Si Tesseract no está instalado o falló el procesamiento.
    """
    if not OCR_AVAILABLE:
        raise ValueError("Tesseract OCR no está instalado. No se pueden procesar imágenes puras.")
    
    try:
        logger.info(f"Procesando imagen con OCR: {file_path}")
        img = Image.open(file_path)
        # Corregir rotación EXIF típica de móviles
        img = ImageOps.exif_transpose(img)
        
        # OSD (Orientation and Script Detection) forzar auto-rotación si no hay EXIF (--psm 1)
        extracted_text = pytesseract.image_to_string(img, lang="spa+eng", config="--psm 1")
        return extracted_text, {}
    except Exception as e:
        raise ValueError(f"Fallo al realizar OCR sobre la imagen: {e}")


def _extract_txt(file_path: str) -> tuple[str, dict]:
    """Lee un archivo de texto plano."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read(), {}


def _extract_csv(file_path: str) -> tuple[str, dict]:
    """
    Transforma un archivo CSV en una representación textual legible de "clave: valor".
    
    Cada fila del CSV se convierte en una línea de texto donde las columnas 
    están precedidas por su nombre de cabecera. Esto facilita la búsqueda 
    semántica al mantener el contexto de cada campo.
    """
    lines = []
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.DictReader(f)
        for row in reader:
            # Convierte cada fila en texto clave:valor
            parts = [f"{key}: {value}" for key, value in row.items() if value]
            lines.append(" | ".join(parts))

    return "\n".join(lines), {}


def _extract_xlsx(file_path: str) -> tuple[str, dict]:
    """
    Lee un XLSX y lo convierte en texto legible.
    Usa los headers como claves, igual que CSV.
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "openpyxl es necesario para leer archivos XLSX. "
            "Instala con: pip install openpyxl"
        )

    wb = load_workbook(file_path, read_only=True, data_only=True)
    all_text = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = list(ws.iter_rows(values_only=True))

        if not rows:
            continue

        # Primera fila como headers
        headers = [str(h) if h else f"col_{i}" for i, h in enumerate(rows[0])]

        for row in rows[1:]:
            parts = []
            for header, value in zip(headers, row):
                if value is not None:
                    parts.append(f"{header}: {value}")
            if parts:
                all_text.append(" | ".join(parts))

    wb.close()
    return "\n".join(all_text), {}


def _extract_docx(file_path: str) -> tuple[str, dict]:
    """Extrae texto de un documento Word (.docx) usando python-docx."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # También extraer texto de tablas
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    meta = {}
    if doc.core_properties.author:
        meta["author"] = doc.core_properties.author
    return "\n".join(paragraphs), meta


def _extract_pptx(file_path: str) -> tuple[str, dict]:
    """Extrae texto de una presentación PowerPoint (.pptx)."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        if parts:
            slides_text.append(f"[Diapositiva {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides_text), {}


def _extract_json(file_path: str) -> tuple[str, dict]:
    """Lee un archivo JSON y lo convierte en texto legible (clave: valor)."""
    import json as json_mod
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        data = json_mod.load(f)

    def flatten(obj, prefix=""):
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                lines.extend(flatten(v, f"{prefix}{k}: "))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                lines.extend(flatten(item, f"{prefix}[{i}] "))
        else:
            lines.append(f"{prefix}{obj}")
        return lines

    return "\n".join(flatten(data)), {}


def _extract_xml(file_path: str) -> tuple[str, dict]:
    """Extrae texto visible de un archivo XML usando BeautifulSoup."""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "lxml-xml")
    return soup.get_text(separator="\n", strip=True), {}


def _extract_html(file_path: str) -> tuple[str, dict]:
    """Extrae texto visible de un archivo HTML usando BeautifulSoup."""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Eliminar scripts y estilos
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    meta = {"html_title": title} if title else {}
    return soup.get_text(separator="\n", strip=True), meta


def clean_text(text: str) -> str:
    """
    Limpia y normaliza el texto crudo extraído de los documentos.

    Este paso es fundamental para eliminar el "ruido" que introducen los 
    diversos extractores (saltos de línea huérfanos, números de página 
    aislados, espacios múltiples) y preparar el texto para una indexación 
    de calidad.

    Operaciones:
        - Colapsa espacios en blanco.
        - Elimina líneas numéricas (típicas de pies de página).
        - Normaliza párrafos colapsando saltos de línea excesivos.
    """
    text = re.sub(r'[^\S\n\t]+', ' ', text)
    text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[ \t]+', '', text, flags=re.MULTILINE)
    return text.strip()


def chunk_text(text: str, size: int = 1000, overlap: int = 200) -> list[str]:
    """
    Divide el texto en fragmentos (chunks) manejables para el LLM.

    Implementa una estrategia de "Semantic Chunking" básico:
        1. Intenta romper por finales de oración (. ! ?) para no cortar 
           ideas a la mitad.
        2. Si no es posible, corta en el último espacio disponible.
        3. Mantiene un solapamiento (overlap) entre fragmentos para preservar 
           la continuidad del contexto semántico en las búsquedas.

    Args:
        text (str): Texto limpio y normalizado.
        size (int): Límite máximo de caracteres por fragmento.
        overlap (int): Cantidad de caracteres que se repiten en el siguiente fragmento.

    Returns:
        list[str]: Lista de fragmentos listos para ser vectorizados.
    """
    if not text:
        return []

    # Separadores de oración (en orden de preferencia)
    sentence_endings = re.compile(r'(?<=[.!?])\s+|(?<=\n)\s*')

    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = start + size

        if end < text_length:
            # Buscar el último final de oración dentro del rango
            segment = text[start:end]
            best_break = -1

            for match in sentence_endings.finditer(segment):
                # Solo considerar quiebres que dejen un chunk razonable (>30% del tamaño)
                if match.start() > size * 0.3:
                    best_break = match.end()

            if best_break > 0:
                end = start + best_break
            else:
                # Fallback: cortar en el último espacio
                last_space = text.rfind(' ', start, end)
                if last_space > start:
                    end = last_space

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - overlap if end < text_length else text_length

    return chunks


STOP_WORDS_ES_EN = {
    "el", "la", "los", "las", "un", "una", "unos", "unas",
    "de", "del", "al", "en", "por", "con", "para", "sin",
    "sobre", "entre", "hacia", "desde", "como", "que", "se",
    "su", "sus", "lo", "es", "son", "no", "si", "ya", "más",
    "muy", "pero", "este", "esta", "estos", "estas", "ese",
    "esa", "esos", "esas", "todo", "toda", "todos", "todas",
    "otro", "otra", "otros", "otras",
    "the", "of", "and", "in", "to", "for", "is", "on",
    "with", "at", "by", "an", "or", "not", "are", "was",
    "be", "has", "had", "its", "from", "this", "that", "which", "also",
}


def normalize_query(query: str) -> str:
    """
    Optimiza la consulta de búsqueda eliminando términos irrelevantes.

    Este pre-procesamiento mejora la precisión de la búsqueda híbrida al 
    centrarse en las palabras con carga semántica, eliminando artículos 
    y preposiciones (Stop Words) comunes en español e inglés.

    Returns:
        str: Consulta normalizada y limpia.
    """
    query = query.strip()
    query = re.sub(r'\s+', ' ', query)
    query = re.sub(r'^[^\w]+|[^\w]+$', '', query)
    # Filtrar stop words para mejorar la búsqueda vectorial/léxica
    words = [w for w in query.split() if w.lower() not in STOP_WORDS_ES_EN]
    return ' '.join(words) if words else query  # fallback si todo son stop words


def deduplicate_chunks(chunks: list[str]) -> list[str]:
    """
    Elimina fragmentos idénticos basándose en su hash SHA256.

    Evita la redundancia en la base de datos vectorial si un documento 
    repite oraciones o párrafos exactos (ej. avisos legales recurrentes).
    """
    seen: set[str] = set()
    unique: list[str] = []

    for chunk in chunks:
        h = hashlib.sha256(chunk.encode('utf-8')).hexdigest()
        if h not in seen:
            seen.add(h)
            unique.append(chunk)

    return unique
